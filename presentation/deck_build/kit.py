from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---- palette (light, conference-safe; mirrors LINDA's light theme) ----
INK    = RGBColor(0x1A,0x20,0x29)
DIM    = RGBColor(0x56,0x60,0x72)
FAINT  = RGBColor(0x8D,0x96,0xA5)
HAIR   = RGBColor(0xD5,0xDB,0xE4)
WELL   = RGBColor(0xF3,0xF5,0xF9)
ACCENT = RGBColor(0x25,0x63,0xEB)
ASOFT  = RGBColor(0xE4,0xEC,0xFD)
RED    = RGBColor(0xCF,0x3A,0x3A)
RSOFT  = RGBColor(0xFB,0xEA,0xEA)
GREEN  = RGBColor(0x1A,0x9D,0x55)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
FONT   = "Arial"

W, H = Inches(10), Inches(5.625)
ML, MR = Inches(0.62), Inches(0.62)
CW = W - ML - MR   # content width

def deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p

def blank(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = WHITE
    return s

def spc(run, hundredths):
    """letter-spacing in 1/100 pt (python-pptx has no API for this)."""
    run.font._rPr.set('spc', str(int(hundredths)))

def tb(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return box, tf

def para(tf, text, size, color, bold=False, first=False, space_before=0,
         line=1.18, align=None, tracking=None, font=FONT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None: p.alignment = align
    p.line_spacing = line
    if space_before: p.space_before = Pt(space_before)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    if tracking: spc(r, tracking)
    return p

def _flat(sh):
    """Kill theme-inherited shadow/effects. shadow.inherit=False alone is not
    enough: add_shape() also emits a <p:style> whose effectRef re-applies one."""
    spPr = sh._element.spPr
    for tag in ('a:effectLst',):
        for e in spPr.findall(qn(tag)): spPr.remove(e)
    spPr.append(spPr.makeelement(qn('a:effectLst'), {}))
    st = sh._element.find(qn('p:style'))
    if st is not None: sh._element.remove(st)
    return sh

def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.shadow.inherit = False
    _flat(sh)
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.text_frame.text = ""
    return sh

def label(slide, x, y, w, text, size=10, color=FAINT, bold=True,
          align=PP_ALIGN.LEFT, tracking=90):
    _, tf = tb(slide, x, y, w, Pt(size*1.6), align=align)
    para(tf, text, size, color, bold=bold, first=True, align=align, tracking=tracking)

def hairline(slide, x, y, w, color=HAIR, thick=Pt(0.75)):
    return rect(slide, x, y, w, thick, fill=color)

# ---------- brand marks ----------
def linda_mark(slide, x, y, scale=1.0, bar_color=ACCENT, text_color=INK,
               wordmark=True, wordmark_rule=364):
    """Redraws live/web/linda-logo.svg: five-bar PSD envelope + baseline + wordmark.
    SVG viewBox is 400x100; `scale` = inches per 100 svg units."""
    u = Inches(scale) / 100.0
    def E(v): return int(round(u * v))
    bars = [(12,42,8,16,0.40),(26,34,8,24,0.62),(40,14,8,44,1.0),
            (54,34,8,24,0.62),(68,46,8,12,0.40)]
    for bx,by,bw,bh,op in bars:
        s = rect(slide, x+E(bx), y+E(by), E(bw), E(bh), fill=bar_color)
        if op < 1.0:
            sf = s.fill.fore_color._xFill.find(qn('a:srgbClr'))
            a = sf.makeelement(qn('a:alpha'), {'val': str(int(op*100000))}); sf.append(a)
    width = wordmark_rule if wordmark else 64
    rect(slide, x+E(12), y+E(66), E(width), E(4), fill=bar_color)
    if wordmark:
        _, tf = tb(slide, x+E(96), y+E(14), E(300), E(52), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, "LINDA", 48*scale*0.72, text_color, bold=True, first=True, tracking=int(700*scale))

def nist_mark(slide):
    """Typographic placeholder, sized + positioned so the official file drops in
    at exactly this box: 1.30in x 0.42in, right-aligned to the margin."""
    x = W - MR - Inches(1.30); y = Inches(0.36)
    _, tf = tb(slide, x, y, Inches(1.30), Inches(0.24), align=PP_ALIGN.RIGHT)
    para(tf, "NIST", 17, INK, bold=True, first=True, align=PP_ALIGN.RIGHT, tracking=180)
    _, tf2 = tb(slide, x, y+Inches(0.255), Inches(1.30), Inches(0.16), align=PP_ALIGN.RIGHT)
    para(tf2, "SURF 2026", 7, FAINT, bold=True, first=True, align=PP_ALIGN.RIGHT, tracking=110)

def footmark(slide):
    linda_mark(slide, ML, H - Inches(0.52), scale=0.30, bar_color=FAINT,
               text_color=FAINT, wordmark_rule=270)

# ---------- standard content slide chrome ----------
def head(slide, kicker, headline, size=27, kicker_color=ACCENT):
    y = Inches(0.36)
    if kicker:
        label(slide, ML, y, Inches(6.0), kicker, size=9.5, color=kicker_color)
        y += Inches(0.30)
    _, tf = tb(slide, ML, y, Inches(7.0), Inches(1.0))
    para(tf, headline, size, INK, bold=True, first=True, line=1.08)
    return y + Inches(0.30 + 0.42*max(1, headline.count("\n")+1))
